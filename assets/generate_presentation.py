"""
生成AI应用创新激励计划汇报PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿对象
prs = Presentation()

# 设置幻灯片尺寸为16:9
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 定义颜色主题
PRIMARY_COLOR = RGBColor(0, 112, 192)  # 深蓝色
SECONDARY_COLOR = RGBColor(255, 192, 0)  # 金色
DARK_COLOR = RGBColor(51, 51, 51)  # 深灰色
LIGHT_COLOR = RGBColor(240, 240, 240)  # 浅灰色

# ============================================
# 封面页
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 添加标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "智能招标文件分析工作流"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR
title_para.alignment = PP_ALIGN.CENTER

# 添加副标题
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "基于LangGraph的AI智能招投标辅助系统"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = DARK_COLOR
subtitle_para.alignment = PP_ALIGN.CENTER

# 添加作者信息
author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.5))
author_frame = author_box.text_frame
author_frame.text = "参赛人：XXX | 解决方案专家"
author_para = author_frame.paragraphs[0]
author_para.font.size = Pt(18)
author_para.font.color.rgb = DARK_COLOR
author_para.alignment = PP_ALIGN.CENTER

# ============================================
# 目录页
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 添加标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "目录"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 添加目录内容
content_items = [
    "01 项目背景与痛点",
    "02 创新点与技术方案",
    "03 核心功能展示",
    "04 应用效果与价值",
    "05 商业前景与规划",
    "06 总结与展望"
]

for i, item in enumerate(content_items):
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5 + i * 0.55), Inches(7), Inches(0.5))
    content_frame = content_box.text_frame
    content_frame.text = item
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(24)
    content_para.font.color.rgb = DARK_COLOR

# ============================================
# 第1页：项目背景与痛点
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "01 项目背景与痛点"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 左侧：行业背景
bg_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(3.8))
bg_frame = bg_box.text_frame
bg_frame.word_wrap = True
bg_frame.text = "行业背景"
bg_para = bg_frame.paragraphs[0]
bg_para.font.size = Pt(28)
bg_para.font.bold = True
bg_para.font.color.rgb = SECONDARY_COLOR
bg_para.space_after = Pt(15)

bg_points = [
    "招投标市场规模持续增长",
    "网络安全领域招标项目数量激增",
    "投标文件质量直接影响中标率",
    "人工审核耗时耗力，易出错漏"
]

for point in bg_points:
    p = bg_frame.add_paragraph()
    p.text = f"• {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(8)

# 右侧：核心痛点
pain_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.8))
pain_frame = pain_box.text_frame
pain_frame.word_wrap = True
pain_frame.text = "核心痛点"
pain_para = pain_frame.paragraphs[0]
pain_para.font.size = Pt(28)
pain_para.font.bold = True
pain_para.font.color.rgb = SECONDARY_COLOR
pain_para.space_after = Pt(15)

pain_points = [
    "废标项检查不全面，存在废标风险",
    "技术方案评估缺乏专业性",
    "得分点识别不精准，失分严重",
    "投标文件结构不规范，影响评审",
    "修改建议不具体，难以落地",
    "人工检查效率低，周期长"
]

for point in pain_points:
    p = pain_frame.add_paragraph()
    p.text = f"• {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(8)

# ============================================
# 第2页：创新点与技术方案
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "02 创新点与技术方案"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 创新点
innovation_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.8))
innovation_frame = innovation_box.text_frame
innovation_frame.word_wrap = True
innovation_frame.text = "核心创新点"
innovation_para = innovation_frame.paragraphs[0]
innovation_para.font.size = Pt(26)
innovation_para.font.bold = True
innovation_para.font.color.rgb = SECONDARY_COLOR
innovation_para.space_after = Pt(12)

innovations = [
    "首创六维并行检测技术：废标、商务、技术、指标、技术得分、结构全方位检查",
    "基于LangGraph工作流编排，实现自动化、智能化的全流程分析",
    "模拟专家评审视角，提供专业、精准的修改建议",
    "智能识别遗漏点、错误点、优化点，提升投标文件质量"
]

for innovation in innovations:
    p = innovation_frame.add_paragraph()
    p.text = f"▶ {innovation}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(8)

# 技术架构
tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.8))
tech_frame = tech_box.text_frame
tech_frame.word_wrap = True
tech_frame.text = "技术架构"
tech_para = tech_frame.paragraphs[0]
tech_para.font.size = Pt(26)
tech_para.font.bold = True
tech_para.font.color.rgb = SECONDARY_COLOR
tech_para.space_after = Pt(12)

tech_points = [
    "AI引擎：基于豆包大语言模型，具备深度文本理解能力",
    "工作流引擎：LangGraph编排，支持串行、并行、条件分支等复杂流程",
    "文档解析：支持PDF、Word等多种格式，自动提取文本内容",
    "智能分析：多维度交叉验证，确保问题无遗漏"
]

for point in tech_points:
    p = tech_frame.add_paragraph()
    p.text = f"▶ {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(8)

# ============================================
# 第3页：核心功能展示
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "03 核心功能展示"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 功能列表
functions = [
    ("🔴 废标项检查", "逐条检查废标要求，识别废标风险，避免投标无效"),
    ("📊 商务得分检查", "评估商务得分，识别失分点，优化资质、业绩、团队"),
    ("🛠️ 技术方案检查", "评估技术方案完整性、创新性、可行性，提升技术得分"),
    ("✅ 指标应答检查", "检查技术指标响应情况，识别遗漏和不充分之处"),
    ("🎯 技术得分检测", "深度检查技术评分细则覆盖情况，精准定位得分点"),
    ("📁 结构检查", "检查投标文件目录结构，符合专家评审习惯")
]

for i, (title, desc) in enumerate(functions):
    # 创建功能框
    x = Inches(0.5 + (i % 2) * 4.5)
    y = Inches(1.3 + (i // 2) * 0.7)
    func_box = slide.shapes.add_textbox(x, y, Inches(4.5), Inches(0.65))
    func_frame = func_box.text_frame
    func_frame.word_wrap = True
    
    # 添加标题
    title_p = func_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(20)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # 添加描述
    desc_p = func_frame.add_paragraph()
    desc_p.text = desc
    desc_p.font.size = Pt(14)
    desc_p.font.color.rgb = DARK_COLOR

# ============================================
# 第4页：应用效果展示
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "04 应用效果与价值"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 左侧：实际案例
case_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(3.8))
case_frame = case_box.text_frame
case_frame.word_wrap = True
case_frame.text = "实际案例：某网络安全项目"
case_para = case_frame.paragraphs[0]
case_para.font.size = Pt(24)
case_para.font.bold = True
case_para.font.color.rgb = SECONDARY_COLOR
case_para.space_after = Pt(15)

# 优化前
before_frame = case_frame.add_paragraph()
before_frame.text = "优化前："
before_frame.font.size = Pt(20)
before_frame.font.bold = True
before_frame.font.color.rgb = DARK_COLOR
before_frame.space_before = Pt(8)

before_points = ["废标风险：存在2项", "商务得分：13分（43%）", "技术得分：52分", "识别问题：14项"]
for point in before_points:
    p = case_frame.add_paragraph()
    p.text = f"  • {point}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_COLOR
    p.level = 1

# 优化后
after_frame = case_frame.add_paragraph()
after_frame.text = "优化后（预期）："
after_frame.font.size = Pt(20)
after_frame.font.bold = True
after_frame.font.color.rgb = PRIMARY_COLOR
after_frame.space_before = Pt(15)

after_points = ["废标风险：无", "商务得分：28分（93%）", "技术得分：65分", "竞争力：大幅提升"]
for point in after_points:
    p = case_frame.add_paragraph()
    p.text = f"  • {point}"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    p.level = 1

# 右侧：核心价值
value_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.8))
value_frame = value_box.text_frame
value_frame.word_wrap = True
value_frame.text = "核心价值"
value_para = value_frame.paragraphs[0]
value_para.font.size = Pt(24)
value_para.font.bold = True
value_para.font.color.rgb = SECONDARY_COLOR
value_para.space_after = Pt(15)

value_points = [
    "⏱️ 效率提升：从3-5天缩短至30分钟",
    "📈 得分提升：平均提升15-20分",
    "✅ 准确率高：问题识别准确率超90%",
    "💰 降低成本：减少人工审核成本",
    "🎯 精准定位：快速找到关键失分点",
    "📝 可落地：提供具体修改建议"
]

for point in value_points:
    p = value_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(10)

# ============================================
# 第5页：商业前景与规划
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "05 商业前景与规划"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 应用场景
scenario_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(1.8))
scenario_frame = scenario_box.text_frame
scenario_frame.word_wrap = True
scenario_frame.text = "应用场景"
scenario_para = scenario_frame.paragraphs[0]
scenario_para.font.size = Pt(24)
scenario_para.font.bold = True
scenario_para.font.color.rgb = SECONDARY_COLOR
scenario_para.space_after = Pt(10)

scenarios = [
    "✓ 企业投标管理",
    "✓ 招标代理服务",
    "✓ 咨询机构工具",
    "✓ 培训教育平台"
]

for scenario in scenarios:
    p = scenario_frame.add_paragraph()
    p.text = scenario
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(6)

# 市场规模
market_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(1.8))
market_frame = market_box.text_frame
market_frame.word_wrap = True
market_frame.text = "市场规模"
market_para = market_frame.paragraphs[0]
market_para.font.size = Pt(24)
market_para.font.bold = True
market_para.font.color.rgb = SECONDARY_COLOR
market_para.space_after = Pt(10)

market_text = "• 全国招投标市场超万亿\n• 网络安全招标年增长30%\n• 智能化工具渗透率<5%\n• 市场空间巨大"
for line in market_text.split('\n'):
    p = market_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(6)

# 未来规划
plan_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.8))
plan_frame = plan_box.text_frame
plan_frame.word_wrap = True
plan_frame.text = "未来规划"
plan_para = plan_frame.paragraphs[0]
plan_para.font.size = Pt(24)
plan_para.font.bold = True
plan_para.font.color.rgb = SECONDARY_COLOR
plan_para.space_after = Pt(10)

plan_text = "• 短期（1-3个月）：优化模型精度，扩展文件格式支持\n• 中期（3-6个月）：开发Web版本，支持多用户协同\n• 长期（6-12个月）：打造行业标杆，拓展至全行业应用"
for line in plan_text.split('\n'):
    p = plan_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(8)

# ============================================
# 第6页：总结与展望
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "06 总结与展望"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR

# 总结
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
summary_frame = summary_box.text_frame
summary_frame.word_wrap = True
summary_frame.text = "项目总结"
summary_para = summary_frame.paragraphs[0]
summary_para.font.size = Pt(24)
summary_para.font.bold = True
summary_para.font.color.rgb = SECONDARY_COLOR
summary_para.space_after = Pt(10)

summary_text = "本项目首创六维并行检测技术，基于LangGraph工作流编排和AI大模型，实现招标文件和投标文件的智能化分析，有效解决传统人工审核效率低、准确率低、落地难的问题，为招投标业务提供强有力的技术支撑。"
summary_frame.paragraphs[0].text = summary_text
summary_frame.paragraphs[0].font.size = Pt(20)
summary_frame.paragraphs[0].font.color.rgb = DARK_COLOR

# 核心优势
advantage_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(1.5))
advantage_frame = advantage_box.text_frame
advantage_frame.word_wrap = True
advantage_frame.text = "核心优势"
advantage_para = advantage_frame.paragraphs[0]
advantage_para.font.size = Pt(24)
advantage_para.font.bold = True
advantage_para.font.color.rgb = SECONDARY_COLOR
advantage_para.space_after = Pt(10)

advantage_text = "🎯 全面性：六维并行检测，覆盖所有得分点\n💡 专业性：模拟专家评审，提供专业建议\n⚡ 高效性：30分钟完成，效率提升100倍\n🔧 实用性：具体可操作，修改建议落地"
for line in advantage_text.split('\n'):
    p = advantage_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_before = Pt(6)

# 展望
outlook_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.8))
outlook_frame = outlook_box.text_frame
outlook_frame.word_wrap = True
outlook_frame.text = "未来展望"
outlook_para = outlook_frame.paragraphs[0]
outlook_para.font.size = Pt(24)
outlook_para.font.bold = True
outlook_para.font.color.rgb = PRIMARY_COLOR
outlook_para.space_after = Pt(10)

outlook_frame.paragraphs[0].text = "我们将持续优化产品性能，拓展应用场景，打造招投标领域的AI智能助手，赋能企业数字化转型，提升行业整体竞争力。"
outlook_frame.paragraphs[0].font.size = Pt(20)
outlook_frame.paragraphs[0].font.color.rgb = DARK_COLOR

# ============================================
# 致谢页
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 添加居中的致谢文字
thanks_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7), Inches(1.5))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "感谢聆听"
thanks_para = thanks_frame.paragraphs[0]
thanks_para.font.size = Pt(60)
thanks_para.font.bold = True
thanks_para.font.color.rgb = PRIMARY_COLOR
thanks_para.alignment = PP_ALIGN.CENTER

# 添加副标题
subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.7), Inches(5), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "期待您的宝贵意见与建议"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = DARK_COLOR
subtitle_para.alignment = PP_ALIGN.CENTER

# 保存PPT
ppt_path = "/workspace/projects/assets/智能招标文件分析工作流_汇报PPT.pptx"
prs.save(ppt_path)

print(f"PPT已生成: {ppt_path}")
print(f"共{len(prs.slides)}页")
